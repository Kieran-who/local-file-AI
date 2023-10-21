from pptx import Presentation
from docx import Document
import openpyxl
from pdfminer.high_level import extract_text as pdfminer_extract
from .segs import segs_processor
from db.docs import new_doc, update_doc_segs, search_docs_path, search_docs_hash, delete_doc
from utils.azure_open_ai import get_chat_completion
import tiktoken
import hashlib
import logging

# Get the logger for the module
logger = logging.getLogger('pdfminer.pdfpage')

# Set the level to ERROR
logger.setLevel(logging.ERROR)


def extract_text(file_path):
    if file_path.endswith('.pdf'):
        text = pdfminer_extract(file_path)

        # Splitting text into lines
        lines = text.split('\n')

        # Joining lines if a line starts with a lowercase letter
        joined_lines = []
        for line in lines:
            if line and line[0].islower():
                if joined_lines:  # ensure joined_lines is not empty
                    joined_lines[-1] = joined_lines[-1] + line
                else:  # if joined_lines is empty, just append the line
                    joined_lines.append(line)
            else:
                joined_lines.append(line)

        # Removing newline characters which might be in the middle of the words
        joined_lines = [line.replace('\n', '') for line in joined_lines]

        # Removing consecutive duplicates
        non_repeat_lines = [joined_lines[i] for i in range(
            len(joined_lines)) if i == 0 or joined_lines[i] != joined_lines[i-1]]

        return '\n'.join(non_repeat_lines)
    elif file_path.endswith('.docx'):
        doc = Document(file_path)
        paragraphs = doc.paragraphs
        non_repeat_paras = [paragraphs[i].text for i in range(
            len(paragraphs)) if i == 0 or paragraphs[i].text != paragraphs[i-1].text]
        return '\n'.join(non_repeat_paras)
    elif file_path.endswith('.pptx'):
        prs = Presentation(file_path)
        text_shapes = [
            shape for slide in prs.slides for shape in slide.shapes if hasattr(shape, "text")]
        non_repeat_shapes = [text_shapes[i].text for i in range(
            len(text_shapes)) if i == 0 or text_shapes[i].text != text_shapes[i-1].text]
        return '\n'.join(non_repeat_shapes)
    elif file_path.endswith('.xlsx'):
        workbook = openpyxl.load_workbook(file_path)
        sheet = workbook.active
        text = ""
        for row in sheet.iter_rows(values_only=True):
            for cell in row:
                text += str(cell) + " "
            text += "\n"
        return text
    elif file_path.endswith('.txt'):
        with open(file_path, 'r') as f:
            text = f.read()
        return text
    else:
        raise ValueError('Unsupported file type')


async def process_doc(doc, pbar):
    # plain function to extract text from doc
    doc_text = extract_text(doc["path"])

    # check if doc exists by checking text hash
    doc_hash = hashlib.sha256(
        f'{doc["path"]}+{doc["name"]}+{doc_text}'.encode()).hexdigest()
    existing_doc = await search_docs_hash(doc_hash)

    if existing_doc:
        pbar.update()
        return ({"id": existing_doc[0]["_additional"]["id"]})
    else:
        # test if old doc exists
        # returns array of docs matching path
        old_doc = await search_docs_path(doc["path"])
        if old_doc:
            for old_doc_item in old_doc:
                if old_doc_item["file_path"] == doc["path"]:

                    await delete_doc(old_doc_item["_additional"]["id"])

        doc_type = doc["path"].split('.')[-1]

        # Tokenize the input string.
        encoding = tiktoken.encoding_for_model("gpt-3.5-turbo")
        num_tokens = len(encoding.encode(doc_text))

        first_5000_tokens = doc_text

        sum_prompt = f"Summarise the below {doc_type} document. It is vital to ensure the summary is as semantically identical to the text but keep it fairly brief. Start your summary with 'This document'"

        # reduce string size for summary generation
        if num_tokens > 5000:
            first_5000_tokens = encoding.encode(doc_text)[:5000]
            first_5000_tokens = encoding.decode(first_5000_tokens)
            sum_prompt = f"Summarise the below {doc_type} document but keep it fairly brief. Due to the document's size, you have been provided only the start of the document. Use this start of the document to write a summary applicable to the whole document. Start your summary with 'This document'"

        # get summary of document
        try:
            text_summary = await get_chat_completion([{"role": "user", "content": f"{sum_prompt}\n###\nDOCUMENT:\n###\n{first_5000_tokens}"}], max_res_tokens=250)
        except Exception as e:
            text_summary = {"message": {
                "content": "There was an error retrieving a summary for this document"}}
        # new doc object
        new_doc_obj = {
            "name": doc["name"],
            "file_path": doc["path"],
            "summary": text_summary["message"]["content"],
            "hash": doc_hash
        }

        # create document
        doc_id = await new_doc(new_doc_obj)

        # process doc and create segments, each with doc id; add each segID to array
        if not doc["path"].endswith('.xlsx'):
            segment_ids = await segs_processor(doc_text, doc_id)

            # add segments to document
            await update_doc_segs(segment_ids, doc_id)
        pbar.update()
        return ({"id": doc_id})
